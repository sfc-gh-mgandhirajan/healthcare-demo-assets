#!/usr/bin/env python3
"""Convert Edwards Clinical Document Intelligence HTML deck content into
Snowflake-branded PPTX using the official January 2026 template."""

import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

TEMPLATE = os.path.expanduser("~/.snowflake/cortex/skills/CoCo_pptx_Skill/snowflake_template.pptx")
OUTPUT = os.path.join(os.path.dirname(os.path.abspath(__file__)), "edwards-lifesciences-deck.pptx")

DK1       = RGBColor(0x26, 0x26, 0x26)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
DK2       = RGBColor(0x11, 0x56, 0x7F)
SF_BLUE   = RGBColor(0x29, 0xB5, 0xE8)
TEAL      = RGBColor(0x71, 0xD3, 0xDC)
ORANGE    = RGBColor(0xFF, 0x9F, 0x36)
VIOLET    = RGBColor(0x7D, 0x44, 0xCF)
BODY_GREY = RGBColor(0x5B, 0x5B, 0x5B)
TBL_GREY  = RGBColor(0x71, 0x71, 0x71)
LIGHT_BG  = RGBColor(0xF5, 0xF5, 0xF5)
BORDER    = RGBColor(0xC8, 0xC8, 0xC8)
ERROR_RED = RGBColor(0xA2, 0x00, 0x00)
GREEN     = RGBColor(0x0A, 0x6E, 0x6E)

FILL_TEXT = {
    "DK2": WHITE, "SF_BLUE": WHITE, "TEAL": DK1, "ORANGE": DK1,
    "VIOLET": WHITE, "GREEN": WHITE, "ERROR_RED": WHITE,
}

# ── Helpers ──

def set_ph(slide, idx, text):
    from pptx.enum.text import MSO_AUTO_SIZE
    from lxml import etree
    ph = slide.placeholders[idx]
    ph.text = text
    ph.text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    ns = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    bodyPr = ph.text_frame._txBody.find(f'{{{ns}}}bodyPr')
    if bodyPr is None:
        bodyPr = etree.SubElement(ph.text_frame._txBody, f'{{{ns}}}bodyPr')
    t_pos = (ph.top or 0) / 914400
    if t_pos < 0.50:
        bodyPr.set('bIns', '0')
    elif 0.60 < t_pos < 1.20:
        bodyPr.set('tIns', '54864')
    if t_pos < 1.20:
        for para in ph.text_frame.paragraphs:
            pPr = para._p.find(f'{{{ns}}}pPr')
            if pPr is None:
                pPr = etree.SubElement(para._p, f'{{{ns}}}pPr')
                para._p.insert(0, pPr)
            pPr.set('indent', '0')
            pPr.set('marL', '0')

def _pad_body_ph(ph):
    from lxml import etree
    t_pos = (ph.top or 0) / 914400
    if t_pos > 1.20:
        ns = 'http://schemas.openxmlformats.org/drawingml/2006/main'
        bodyPr = ph.text_frame._txBody.find(f'{{{ns}}}bodyPr')
        if bodyPr is None:
            bodyPr = etree.SubElement(ph.text_frame._txBody, f'{{{ns}}}bodyPr')
        bodyPr.set('bIns', '91440')

def set_ph_lines(slide, idx, lines, font_size=None):
    ph = slide.placeholders[idx]
    tf = ph.text_frame
    tf.clear()
    _pad_body_ph(ph)
    lines = [l for l in lines if l.strip()]
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        if font_size:
            p.font.size = Pt(font_size)

def set_ph_sections(slide, idx, sections, heading_size=None, body_size=None):
    from lxml import etree
    ph = slide.placeholders[idx]
    tf = ph.text_frame
    tf.clear()
    _pad_body_ph(ph)
    first = True
    for heading, body_lines in sections:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        p.level = 0
        if not first:
            ns = 'http://schemas.openxmlformats.org/drawingml/2006/main'
            pPr = p._p.find(f'{{{ns}}}pPr')
            if pPr is None:
                pPr = etree.SubElement(p._p, f'{{{ns}}}pPr')
                p._p.insert(0, pPr)
            spcBef = etree.SubElement(pPr, f'{{{ns}}}spcBef')
            spcPts = etree.SubElement(spcBef, f'{{{ns}}}spcPts')
            spcPts.set('val', '1400')
        first = False
        run = p.add_run()
        run.text = heading
        run.font.bold = True
        run.font.color.rgb = DK2
        if heading_size:
            run.font.size = Pt(heading_size)
        for line in body_lines:
            bp = tf.add_paragraph()
            bp.level = 1
            bp.text = line
            if body_size:
                bp.font.size = Pt(body_size)

def set_ph_bold_keywords(slide, idx, items, kw_size=None, body_size=None):
    ph = slide.placeholders[idx]
    tf = ph.text_frame
    tf.clear()
    _pad_body_ph(ph)
    for i, (keyword, description) in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        r1 = p.add_run()
        r1.text = f"{keyword}: "
        r1.font.bold = True
        if kw_size: r1.font.size = Pt(kw_size)
        r2 = p.add_run()
        r2.text = description
        if body_size: r2.font.size = Pt(body_size)

def add_shape_text(slide, shape_type, left, top, width, height,
                   text, fill_colour, font_colour,
                   font_size=10, bold=False, alignment=PP_ALIGN.CENTER):
    shape = slide.shapes.add_shape(
        shape_type,
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_colour
    shape.line.fill.background()
    if width <= 2.0 and '\n' not in text and ' ' in text:
        text = text.replace(' ', '\n')
    tf = shape.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Pt(4)
    tf.margin_right = Pt(4)
    tf.margin_top = Pt(2)
    tf.margin_bottom = Pt(2)
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = "Arial"
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = font_colour
    p.alignment = alignment
    return shape

def add_connector_arrow(slide, x1, y1, x2, y2):
    from pptx.oxml.ns import qn
    from lxml import etree
    cx = slide.shapes.add_shape(
        MSO_SHAPE.RIGHT_ARROW, Inches(x1), Inches(y1),
        Inches(x2 - x1), Inches(0.25))
    cx.fill.solid()
    cx.fill.fore_color.rgb = SF_BLUE
    cx.line.fill.background()
    return cx

# ── Build Deck ──

prs = Presentation(TEMPLATE)

# Remove sample slides (keep only layouts)
while len(prs.slides) > 0:
    rId = prs.slides._sldIdLst[0].get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id')
    prs.part.drop_rel(rId)
    prs.slides._sldIdLst.remove(prs.slides._sldIdLst[0])

# ═══════════════════════════════════════════
# SLIDE 1: COVER (Layout 13)
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[13])
set_ph(slide, 3, "CLINICAL DOCUMENT\nINTELLIGENCE")
set_ph(slide, 0, "Edwards Lifesciences × Snowflake")
set_ph(slide, 2, "Structural Heart Trials  |  TMF Governance Demo")

# ═══════════════════════════════════════════
# SLIDE 2: SCENARIO (Layout 5 — 1 col)
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[5])
set_ph(slide, 0, "WHEN AN INSPECTION IS ANNOUNCED")
set_ph(slide, 2, "A typical clinical trial lifecycle moment — everyone has questions")
set_ph_sections(slide, 1, [
    ("The Scenario", [
        "A medical device or pharma company is running multiple pivotal trials — each generating hundreds of clinical documents in an eTMF system",
        "A regulatory agency announces a pre-approval inspection — the organization has weeks, not months, to demonstrate audit readiness",
        "Protocols have been amended multiple times; informed consent forms vary by site — version control is fragmented",
        "Leadership needs a competitive landscape briefing — but trial registry data is scattered and stale",
        "Every persona has an urgent, unanswered question — and today, answering any of them takes days of manual effort",
    ]),
], heading_size=12, body_size=10)

# ═══════════════════════════════════════════
# SLIDE 4: FOUR PERSONAS (Layout 8 — 4 col)
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[8])
set_ph(slide, 0, "FOUR PERSONAS, FOUR QUESTIONS")
set_ph(slide, 5, "Each persona has a critical question that takes days to answer manually")
set_ph_sections(slide, 1, [
    ("CMO / VP Clinical Ops", [
        "\"Which trials have the lowest document completeness — and where is our risk exposure highest?\"",
    ]),
], heading_size=10, body_size=9)
set_ph_sections(slide, 2, [
    ("TMF Manager", [
        "\"Are all sites on the current protocol version? Are any consent forms outdated?\"",
    ]),
], heading_size=10, body_size=9)
set_ph_sections(slide, 3, [
    ("VP Regulatory Affairs", [
        "\"Can I get an audit readiness report — completeness by TMF zone, gaps, and action items?\"",
    ]),
], heading_size=10, body_size=9)
set_ph_sections(slide, 4, [
    ("Head of Clinical Strategy", [
        "\"What are competitors doing in our therapeutic area? How do their trial designs compare?\"",
    ]),
], heading_size=10, body_size=9)

# ═══════════════════════════════════════════
# SLIDE 5: CHAPTER — THE IMPERATIVE (Layout 18)
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[18])
set_ph(slide, 1, "THE IMPERATIVE")

# ═══════════════════════════════════════════
# SLIDE 6: WHY NOW — TMF GOVERNANCE GAP (Layout 6 — 2 col)
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_ph(slide, 0, "THE INDUSTRY-WIDE TMF GOVERNANCE GAP")
set_ph(slide, 3, "Two pain points are converging — and manual processes cannot keep pace")
set_ph_sections(slide, 1, [
    ("TMF Audit Readiness Takes Days", [
        "Inspection announced — TMF teams manually pull document inventories trial by trial, zone by zone",
        "Version control gaps — protocol amendments and informed consent forms drift across sites",
        "No cross-trial visibility — governance boards rely on manual spreadsheet consolidation",
        "Reactive, not proactive — gaps are found during inspections instead of being prevented",
    ]),
], heading_size=11, body_size=9)
set_ph_sections(slide, 2, [
    ("Competitive Intelligence Is Manual", [
        "Board meeting prep — strategy teams manually search ClinicalTrials.gov and copy results into decks",
        "Data is stale by the time it reaches leadership — competitor trials change status weekly",
        "No cross-reference between internal protocol designs and competitor trial registrations",
        "Multiple concurrent trials — monitoring competitive landscapes manually is unsustainable",
    ]),
], heading_size=11, body_size=9)

# ═══════════════════════════════════════════
# SLIDE 7: CHAPTER — THE ARCHITECTURE (Layout 18)
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[18])
set_ph(slide, 1, "THE\nARCHITECTURE")

# ═══════════════════════════════════════════
# SLIDE 8: ARCHITECTURE — END TO END (Layout 0 — free canvas)
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[0])
set_ph(slide, 0, "END-TO-END ON SNOWFLAKE")
set_ph(slide, 1, "From eTMF system to governed AI console — every layer runs inside your Snowflake account")

# Row 1: Ingestion pipeline
y1 = 1.50
boxes_r1 = [
    ("Veeva Vault\nOpenFlow", DK2),
    ("@ETMF_STAGE\nClinical PDFs", SF_BLUE),
    ("AI_PARSE +\nAI_EXTRACT", TEAL),
    ("PARSED_DOCS\n+ METADATA", DK2),
]
for i, (label, color) in enumerate(boxes_r1):
    x = 0.50 + i * 2.40
    txt_color = WHITE if color != TEAL else DK1
    add_shape_text(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, y1, 1.80, 0.80,
                   label, color, txt_color, font_size=9, bold=True)
    if i < len(boxes_r1) - 1:
        add_shape_text(slide, MSO_SHAPE.RIGHT_ARROW, x + 1.90, y1 + 0.25, 0.40, 0.30,
                       "", LIGHT_BG, DK1, font_size=8)

# Row 2: Intelligence layer
y2 = 2.70
boxes_r2 = [
    ("Cortex\nSearch", SF_BLUE),
    ("Semantic\nView", DK2),
    ("Audit Report\nGen", GREEN),
    ("ClinicalTrials\nCKE", ORANGE),
]
for i, (label, color) in enumerate(boxes_r2):
    x = 0.50 + i * 2.40
    txt_color = WHITE if color not in (TEAL, ORANGE) else DK1
    add_shape_text(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, y2, 1.80, 0.80,
                   label, color, txt_color, font_size=9, bold=True)

# Row 3: Agent + App
y3 = 3.90
add_shape_text(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 2.00, y3, 2.50, 0.80,
               "CORTEX AGENT\n4-Tool Orchestration", DK2, WHITE, font_size=10, bold=True)
add_shape_text(slide, MSO_SHAPE.RIGHT_ARROW, 4.60, y3 + 0.25, 0.40, 0.30,
               "", LIGHT_BG, DK1, font_size=8)
add_shape_text(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 5.10, y3, 2.50, 0.80,
               "REACT + FLASK\nSPCS App", SF_BLUE, WHITE, font_size=10, bold=True)

# ═══════════════════════════════════════════
# SLIDE 9: CHAPTER — APP WALKTHROUGH (Layout 18)
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[18])
set_ph(slide, 1, "APP\nWALKTHROUGH")

# ═══════════════════════════════════════════
# SLIDE 10: SCREEN 1 — TRIAL PORTFOLIO (Layout 7 — 3 col)
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[7])
set_ph(slide, 0, "SCREEN 1: TRIAL PORTFOLIO OVERVIEW")
set_ph(slide, 4, "Persona: CMO / VP Clinical Operations")
set_ph_sections(slide, 1, [
    ("Completeness Scoring", [
        "Each trial card shows real-time document completeness",
        "Gaps are immediately visible without opening a spreadsheet",
    ]),
], heading_size=11, body_size=9)
set_ph_sections(slide, 2, [
    ("RAG Status Badges", [
        "Red/Amber/Green badges per trial based on doc age and completeness",
        "CMO sees risk at a glance — no meeting required",
    ]),
], heading_size=11, body_size=9)
set_ph_sections(slide, 3, [
    ("Document Freshness", [
        "Average days since last document upload per trial",
        "Flags trials needing attention before inspection",
    ]),
], heading_size=11, body_size=9)

# ═══════════════════════════════════════════
# SLIDE 11: SCREEN 2 — TRIAL DRILLDOWN (Layout 7 — 3 col)
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[7])
set_ph(slide, 0, "SCREEN 2: TRIAL DRILLDOWN — DOCUMENT INVENTORY")
set_ph(slide, 4, "Persona: TMF Manager / Document Controller")
set_ph_sections(slide, 1, [
    ("Document Inventory", [
        "Click into a trial and see every document — protocols, ICFs, regulatory submissions",
        "Full document lineage visible per trial",
    ]),
], heading_size=11, body_size=9)
set_ph_sections(slide, 2, [
    ("Version Tracking", [
        "Amendment history tracked via AI_EXTRACT",
        "Identifies sites using outdated consent forms",
    ]),
], heading_size=11, body_size=9)
set_ph_sections(slide, 3, [
    ("Functional Group Gaps", [
        "Filter by functional group (Clinical Ops, Data Mgmt, Safety)",
        "See which teams have missing or outdated uploads",
    ]),
], heading_size=11, body_size=9)

# ═══════════════════════════════════════════
# SLIDE 12: SCREEN 3 — AI AGENT (Layout 5 — 1 col)
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[5])
set_ph(slide, 0, "SCREEN 3: AI AGENT — THE UNIFIER")
set_ph(slide, 2, "Persona: All personas — one agent, four tools")
set_ph_sections(slide, 1, [
    ("Cortex Search (etmf_search)", [
        "Semantic search across parsed clinical documents",
        "Grounded answers with source citations",
    ]),
    ("Cortex Analyst (tmf_analytics)", [
        "Text-to-SQL over the Semantic View",
        "SQL results rendered as inline tables",
    ]),
    ("Audit Report Gen (audit_report_gen)", [
        "Custom stored procedure calls CORTEX.COMPLETE",
        "Generates formatted audit readiness reports on demand",
    ]),
    ("ClinicalTrials.gov (clinical_trials_registry)", [
        "5.7M trial records via Snowflake Marketplace CKE",
        "Competitive intelligence in seconds",
    ]),
], heading_size=11, body_size=9)

# ═══════════════════════════════════════════
# SLIDE 13: SCREEN 4 — AUDIT REPORTS (Layout 7 — 3 col)
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[7])
set_ph(slide, 0, "SCREEN 4: ONE-CLICK AUDIT REPORTS")
set_ph(slide, 4, "Persona: VP Regulatory Affairs / Quality")
set_ph_sections(slide, 1, [
    ("Executive Summary", [
        "One-click per trial produces comprehensive summary",
        "Overall completeness score and RAG assessment",
    ]),
], heading_size=11, body_size=9)
set_ph_sections(slide, 2, [
    ("Zone-by-Zone Analysis", [
        "ICH-GCP mandated TMF zones analyzed individually",
        "Specific gaps identified per zone",
    ]),
], heading_size=11, body_size=9)
set_ph_sections(slide, 3, [
    ("Risk Items & Actions", [
        "Missing documents, upcoming IRB expirations",
        "Prioritized action items generated on demand",
    ]),
], heading_size=11, body_size=9)

# ═══════════════════════════════════════════
# SLIDE 14: COMPETITIVE INTELLIGENCE (Layout 5 — 1 col)
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[5])
set_ph(slide, 0, "COMPETITIVE INTELLIGENCE — CLINICALTRIALS.GOV")
set_ph(slide, 2, "Persona: Head of Clinical Strategy — the showstopper")
set_ph_sections(slide, 1, [
    ("Competitor Landscape", [
        "Ask about competitor trials and get NCT numbers, designs, enrollment targets, and status in seconds",
    ]),
    ("Internal + External Cross-Reference", [
        "Cross-reference internal protocol designs with competitor trial registrations",
        "Compare enrollment criteria against similar trials globally",
    ]),
    ("Always-Current Data", [
        "5.7M trial records from ClinicalTrials.gov via Snowflake Marketplace",
        "Always current — no more stale slide decks",
    ]),
], heading_size=11, body_size=10)

# ═══════════════════════════════════════════
# SLIDE 15: CHAPTER — UNDER THE HOOD (Layout 18)
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[18])
set_ph(slide, 1, "UNDER\nTHE HOOD")

# ═══════════════════════════════════════════
# SLIDE 16: SNOWFLAKE FEATURES TABLE (Layout 0 — table)
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[0])
set_ph(slide, 0, "SNOWFLAKE FEATURES AT EVERY LAYER")
set_ph(slide, 1, "The complete stack — zero structured data in the live demo, 100% unstructured document AI")

headers = ["Layer", "Snowflake Feature", "What It Does"]
data = [
    ["Ingestion", "OpenFlow Connector + Stage", "Lands eTMF docs from Veeva Vault into @ETMF_STAGE"],
    ["Extraction", "AI_PARSE_DOCUMENT + AI_EXTRACT", "Full-text OCR + structured metadata extraction"],
    ["Search", "Cortex Search Service", "Semantic search across parsed document corpus"],
    ["Analytics", "Semantic View (Cortex Analyst)", "7 metrics, 14 dimensions over DOCUMENT_METADATA"],
    ["Orchestration", "Cortex Agent (4 tools)", "Auto-routing between search, analyst, audit gen, registry"],
    ["Intelligence", "ClinicalTrials.gov CKE", "5.7M trial records for competitive intelligence"],
    ["Application", "SPCS (React + Flask)", "Full-stack app deployed inside Snowflake"],
]
n_rows = len(data) + 1
n_cols = len(headers)
tbl_shape = slide.shapes.add_table(
    n_rows, n_cols,
    Inches(0.40), Inches(1.40), Inches(9.10), Inches(0.40 * n_rows))
tbl = tbl_shape.table

for ci, h in enumerate(headers):
    cell = tbl.cell(0, ci)
    cell.text = h
    cell.fill.solid(); cell.fill.fore_color.rgb = DK2
    for p in cell.text_frame.paragraphs:
        p.font.size = Pt(11); p.font.bold = True
        p.font.color.rgb = WHITE; p.font.name = "Arial"

col_widths = [Inches(1.40), Inches(3.50), Inches(4.20)]
for ci, cw in enumerate(col_widths):
    tbl.columns[ci].width = cw

for ri, row in enumerate(data):
    for ci, val in enumerate(row):
        cell = tbl.cell(ri + 1, ci)
        cell.text = str(val)
        cell.fill.solid()
        cell.fill.fore_color.rgb = WHITE if ri % 2 == 0 else LIGHT_BG
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(10)
            p.font.color.rgb = TBL_GREY
            p.font.name = "Arial"
            if ci == 0:
                p.font.bold = True
                p.font.color.rgb = DK2

# ═══════════════════════════════════════════
# SLIDE 17: CHAPTER — SUMMARY (Layout 18)
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[18])
set_ph(slide, 1, "SUMMARY &\nPATH FORWARD")

# ═══════════════════════════════════════════
# SLIDE 18: PERSONA VALUE SUMMARY (Layout 7 — 3 col)
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[7])
set_ph(slide, 0, "ONE PLATFORM. NOT JUST AN EDW.")
set_ph(slide, 4, "The platform is a clinical intelligence platform — ready to transform TMF governance")
set_ph_sections(slide, 1, [
    ("For the TMF Manager", [
        "3-day audit prep scramble → 30-second query",
        "Real-time document completeness and version tracking",
        "Functional group accountability",
    ]),
], heading_size=11, body_size=9)
set_ph_sections(slide, 2, [
    ("For VP Regulatory", [
        "Real-time governance visibility across every trial, zone, and jurisdiction",
        "One-click audit reports with prioritized action items",
    ]),
], heading_size=11, body_size=9)
set_ph_sections(slide, 3, [
    ("For Clinical Strategy", [
        "Competitive intelligence at the speed of a question",
        "5.7M trial records — internal + external in a single query",
        "Always current",
    ]),
], heading_size=11, body_size=9)

# ═══════════════════════════════════════════
# SLIDE 19: PATH FORWARD — CHEVRON (Layout 0)
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[0])
set_ph(slide, 0, "PATH FORWARD — 4-6 WEEK PILOT")
set_ph(slide, 1, "We can help you get there — let's talk about what a pilot looks like for your team")

steps = [
    ("1. CONNECT", "Land Veeva Vault eTMF\ndocs via OpenFlow", DK2),
    ("2. EXTRACT", "AI_PARSE + AI_EXTRACT\nfor searchable corpus", SF_BLUE),
    ("3. GOVERN", "Cortex Agent orchestrates\nsearch, analytics, reports", GREEN),
    ("4. DEPLOY", "SPCS app gives every\npersona a console", TEAL),
]
for i, (title, desc, color) in enumerate(steps):
    x = 0.40 + i * 2.35
    txt_color = WHITE if color not in (TEAL, ORANGE) else DK1
    add_shape_text(slide, MSO_SHAPE.CHEVRON, x, 1.60, 2.20, 0.90,
                   title, color, txt_color, font_size=10, bold=True)
    add_shape_text(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x + 0.10, 2.70, 2.00, 0.90,
                   desc, LIGHT_BG, DK1, font_size=9, bold=False)

# Quick wins vs strategic outcomes
add_shape_text(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 0.40, 3.90, 4.30, 1.10,
               "QUICK WINS (WEEKS 1-2)\nParse documents with AI_PARSE_DOCUMENT\nExtract metadata with AI_EXTRACT\nDeploy Cortex Search + Cortex Analyst",
               WHITE, DK1, font_size=9, bold=False, alignment=PP_ALIGN.LEFT)
add_shape_text(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 5.10, 3.90, 4.30, 1.10,
               "STRATEGIC OUTCOMES (WEEKS 3-6)\nCortex Agent with 4-tool orchestration\nClinicalTrials.gov CKE for competitive intel\nFull SPCS app with role-based access",
               DK2, WHITE, font_size=9, bold=False, alignment=PP_ALIGN.LEFT)

# ═══════════════════════════════════════════
# SLIDE 20: THANK YOU (Layout 28)
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[28])
set_ph(slide, 1, "THANK\nYOU")

# ── Save ──
prs.save(OUTPUT)
print(f"PPTX generated: {OUTPUT}")
print(f"File size: {os.path.getsize(OUTPUT) / 1024:.0f} KB")
print(f"Slides: {len(prs.slides)}")
